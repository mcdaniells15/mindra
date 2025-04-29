import os
from typing import List, Dict
import google.generativeai as genai
from dotenv import load_dotenv
from PyPDF2 import PdfReader
import ast
from langchain.text_splitter import RecursiveCharacterTextSplitter
from config import Config
import asyncio
from concurrent.futures import ThreadPoolExecutor
from pptx import Presentation
from docx import Document

# Load environment variables and configure Gemini API
load_dotenv()
genai.configure(api_key=Config.API_KEY)
model = genai.GenerativeModel(Config.API_MODEL)

# Create a thread pool for concurrent processing
executor = ThreadPoolExecutor(max_workers=Config.MAX_CONCURRENT_REQUESTS)

def extract_text_from_pdf(pdf_path: str) -> str:
    """Extract text content from a PDF file."""
    pdf_reader = PdfReader(pdf_path)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_powerpoint(pptx_path: str) -> str:
    """Extract text content from a PowerPoint file."""
    prs = Presentation(pptx_path)
    text = ""
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    
    return text

def extract_text_from_word(doc_path: str) -> str:
    """Extract text content from a Word document."""
    doc = Document(doc_path)
    text = ""
    
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    
    # Extract text from tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text += cell.text + "\n"
    
    return text

def extract_text_from_file(file_path: str) -> str:
    """Extract text content from a file based on its extension."""
    file_extension = os.path.splitext(file_path)[1].lower()
    
    if file_extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension in ['.pptx', '.ppt']:
        return extract_text_from_powerpoint(file_path)
    elif file_extension in ['.docx', '.doc']:
        return extract_text_from_word(file_path)
    else:
        raise ValueError(f"Unsupported file type: {file_extension}")

def split_text(text: str) -> List[str]:
    """Split text into smaller chunks for processing using LangChain's text splitter."""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=Config.CHUNK_SIZE,
        chunk_overlap=Config.CHUNK_OVERLAP,
        length_function=len
    )
    return text_splitter.split_text(text)

async def generate_summary(text: str) -> str:
    """Generate a comprehensive summary of the text using Gemini API."""
    loop = asyncio.get_event_loop()
    prompt = f"""Please provide a detailed and structured summary of the following text. 
    Follow this format:

    1. MAIN POINTS:
    - List the 3-5 most important points
    - Explain each point briefly but clearly

    2. KEY CONCEPTS:
    - Identify and explain the main concepts discussed
    - Show how they relate to each other

    3. DETAILED EXPLANATION:
    - Provide a thorough explanation of the content
    - Include relevant examples or evidence
    - Explain the significance of each point

    4. CONCLUSION:
    - Summarize the overall message
    - Highlight the most important takeaways
    - Explain the practical implications

    Text to summarize:
    {text}
    """
    response = await loop.run_in_executor(executor, model.generate_content, prompt)
    return response.text

async def generate_quiz(text: str, num_questions: int = 5, difficulty: str = "intermediate") -> List[Dict]:
    """Generate quiz questions based on the text and difficulty level."""
    difficulty_prompt = Config.get_difficulty_prompt(difficulty)
    
    prompt = f"""Generate a {num_questions}-question multiple-choice quiz based on this text.
    {difficulty_prompt}
    
    For each question include:
    1. A clear question
    2. Four options labeled A) through D)
    3. The correct answer letter
    4. A brief explanation

    Text to base the quiz on:
    {text}

    Respond in this exact format:
    [
      {{
        "question": "Write the question here?",
        "options": [
          "A) First option",
          "B) Second option",
          "C) Third option",
          "D) Fourth option"
        ],
        "correct_answer": "A",
        "explanation": "Explanation why this is correct"
      }}
    ]"""

    try:
        loop = asyncio.get_event_loop()
        response = await loop.run_in_executor(executor, model.generate_content, prompt)
        result = response.text.strip()
        
        # Clean up response and ensure proper format
        result = result.replace('```python', '').replace('```json', '').replace('```', '').strip()
        
        if not (result.startswith('[') and result.endswith(']')):
            raise ValueError("Response is not in the correct list format")

        # Parse the response safely
        try:
            questions = ast.literal_eval(result)
        except:
            local_dict = {}
            exec(f"quiz = {result}", {}, local_dict)
            questions = local_dict['quiz']
        
        # Validate response structure
        if not isinstance(questions, list):
            raise ValueError("Response is not a list")
        
        for q in questions:
            if not isinstance(q, dict):
                raise ValueError("Question is not a dictionary")
            if not all(key in q for key in ['question', 'options', 'correct_answer', 'explanation']):
                raise ValueError("Question missing required fields")
            if not isinstance(q['options'], list) or len(q['options']) != 4:
                raise ValueError("Question must have exactly 4 options")
        
        return questions

    except Exception as e:
        print(f"Quiz generation error: {str(e)}")
        print(f"Raw response: {result}")
        return [{
            "question": "Failed to generate quiz questions",
            "options": [
                "A) Try generating again",
                "B) Check the text content",
                "C) Ensure text is clear",
                "D) Contact support"
            ],
            "correct_answer": "A",
            "explanation": "The AI had trouble generating questions. Please try again."
        }]

def format_quiz_for_display(quiz_questions: List[Dict]) -> str:
    """Format quiz questions for display in a readable format."""
    try:
        formatted_quiz = ""
        for i, q in enumerate(quiz_questions, 1):
            formatted_quiz += f"\nQuestion {i}:\n{q.get('question', 'Missing question')}\n\n"
            for option in q.get('options', ['No options available']):
                formatted_quiz += f"{option}\n"
            formatted_quiz += f"\nCorrect Answer: {q.get('correct_answer', 'Not available')}\n"
            formatted_quiz += f"Explanation: {q.get('explanation', 'No explanation available')}\n"
            formatted_quiz += "-" * 50 + "\n"
        return formatted_quiz
    except Exception as e:
        return f"Error formatting quiz: {str(e)}\nPlease try generating the quiz again."

async def generate_practical_explanation(text: str, difficulty: str = "intermediate") -> str:
    """Generate practical explanations and examples based on the document content and difficulty level."""
    difficulty_prompt = Config.get_difficulty_prompt(difficulty)
    
    prompt = f"""Please provide practical explanations and examples based on the following document content.
    {difficulty_prompt}
    
    Focus on:
    1. Practical applications of the concepts
    2. Real-world examples
    3. Step-by-step explanations where applicable
    4. Common use cases
    
    Document content:
    {text}
    """
    loop = asyncio.get_event_loop()
    response = await loop.run_in_executor(executor, model.generate_content, prompt)
    return response.text

async def chat_with_document(text: str, question: str, difficulty: str = "intermediate") -> str:
    """Generate a response to a question about the document content with appropriate difficulty level."""
    difficulty_prompt = Config.get_difficulty_prompt(difficulty)
    
    prompt = f"""You are a helpful assistant that answers questions about the following document content.
    {difficulty_prompt}
    Please provide a clear and accurate response to the question based on the information in the document.
    If the information is not available in the document, say so clearly.
    
    Document content:
    {text}
    
    Question: {question}
    """
    loop = asyncio.get_event_loop()
    response = await loop.run_in_executor(executor, model.generate_content, prompt)
    return response.text 